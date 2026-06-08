# -*- coding: utf-8 -*-
"""第2页重排为嵌套式：外层=FQL主动决策闭环(细化)，内层=LIO-SLAM被动建图。
保留 slide1；复用 slide2 的 3 张照片(机器人/IMU/点云)重新定位；删除其余旧形状后重建。
另存 研究内容三_总体框架图_嵌套版.pptx，不改原文件。"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

SRC = "研究内容三_总体框架图.pptx"
DST = "研究内容三_总体框架图_嵌套版.pptx"

NAVY="0E2841"; TEAL="1C7293"; TEAL_D="0F4C5C"; GREY="5B6B7C"
AMBER="E9A23B"; AMBER_BAR="C77F1A"; AMBER_TXT="8F5E10"
OUT_FILL="FFFBF3"; IN_FILL="EAF3F5"; SUB_FILL="FBE7C6"; WHITE="FFFFFF"
YH="Microsoft YaHei"; KT="楷体"

prs = Presentation(SRC)
slide = prs.slides[1]

def rgb(h): return RGBColor.from_string(h)
def set_cjk(run, name):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin","a:ea","a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", name)

def box(x,y,w,h,fill,line,dash=False,lw=1.0,radius=0.06):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x),Inches(y),Inches(w),Inches(h))
    try: sp.adjustments[0]=radius
    except Exception: pass
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=rgb(fill)
    if line is None: sp.line.fill.background()
    else:
        sp.line.color.rgb=rgb(line); sp.line.width=Pt(lw)
        if dash:
            ln=sp.line._get_or_add_ln(); d=ln.makeelement(qn('a:prstDash'),{}); d.set('val','dash'); ln.append(d)
    sp.shadow.inherit=False
    return sp

def text(x,y,w,h,runs,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,wrap=True):
    """runs: list of (text,size,color,font,bold[,newpara])"""
    tb=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=wrap
    tf.margin_left=tf.margin_right=Pt(3); tf.margin_top=Pt(2); tf.margin_bottom=Pt(2)
    tf.vertical_anchor=anchor
    p=tf.paragraphs[0]; p.alignment=align
    for i,r in enumerate(runs):
        t,sz,col,fn,bd = r[0],r[1],r[2],r[3],r[4]
        newpara = len(r)>5 and r[5]
        if newpara:
            p=tf.add_paragraph(); p.alignment=align
        run=p.add_run(); run.text=t
        run.font.size=Pt(sz); run.font.bold=bd; run.font.color.rgb=rgb(col); set_cjk(run,fn)
    return tb

def card(x,y,w,h,title,sub,line,tcolor,fill=WHITE,tsz=11,ssz=9):
    box(x,y,w,h,fill,line,dash=False,lw=1.0,radius=0.08)
    text(x,y,w,h,[(title,tsz,tcolor,YH,True),(sub,ssz,GREY,KT,False,True)])

def barrow(x,y,w,h,shape=MSO_SHAPE.RIGHT_ARROW,fill=NAVY):
    sp=slide.shapes.add_shape(shape,Inches(x),Inches(y),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=rgb(fill); sp.line.fill.background(); sp.shadow.inherit=False
    return sp

def dline(x1,y1,x2,y2,color=AMBER,w=2.0,dash=True):
    cn=slide.shapes.add_connector(2,Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    cn.line.color.rgb=rgb(color); cn.line.width=Pt(w)
    if dash:
        ln=cn.line._get_or_add_ln(); d=ln.makeelement(qn('a:prstDash'),{}); d.set('val','dash'); ln.append(d)
    cn.shadow.inherit=False
    return cn

# ---------- 复用照片：找出 3 张图片并重定位 ----------
pics=[s for s in slide.shapes if s.shape_type==13]
pics.sort(key=lambda s:(s.width*s.height), reverse=True)  # robot > IMU > pointcloud
robot = pics[0] if len(pics)>0 else None
imu   = pics[1] if len(pics)>1 else None
pc    = pics[2] if len(pics)>2 else None

# 删除其余形状
for s in [sh for sh in slide.shapes if sh.shape_type!=13]:
    s._element.getparent().remove(s._element)

def place(p,x,y,w,h):
    if p is None: return
    p.left=Inches(x); p.top=Inches(y); p.width=Inches(w); p.height=Inches(h)

# ===================== 标题 =====================
text(0.4,0.14,12.4,0.5,[("图2  总体方法框架：退化感知驱动的具身主动 SLAM 闭环",18,NAVY,YH,True)],align=PP_ALIGN.LEFT)

# ===================== 外层框：FQL 主动决策闭环 =====================
OX,OY,OW,OH = 0.3,0.72,12.73,6.62
box(OX,OY,OW,OH,OUT_FILL,AMBER,dash=True,lw=1.5,radius=0.035)
box(OX,OY,OW,0.42,AMBER_BAR,None,radius=0.03)
text(OX+0.2,OY,OW-0.4,0.42,[("主动感知与自适应导航 · FQL 上层决策闭环（外层）",14,WHITE,YH,True)],align=PP_ALIGN.LEFT)

# ===================== FQL 决策行（顶部，细化）=====================
RY,RH = 1.34,2.0
# 模块 x
m1x,m1w=0.55,1.9
m2x,m2w=2.57,2.1
m3x,m3w=4.79,3.7
m4x,m4w=8.61,2.25
m5x,m5w=10.98,1.7

card(m1x,RY,m1w,RH,"候选生成","前沿点 + 潜在回环点\n构成动作空间 𝒜ₜ",AMBER,AMBER_TXT)
card(m2x,RY,m2w,RH,"状态编码 sₜ","地图熵 H(m)·协方差 Σ\n可定位性 η·局部几何",AMBER,AMBER_TXT)

# M3 FQL 双策略核心（细化：3 子框）
box(m3x,RY,m3w,RH,WHITE,AMBER,dash=False,lw=1.5,radius=0.05)
text(m3x,RY+0.05,m3w,0.3,[("FQL 双策略核心（无 BPTT · 单步推理）",10.5,AMBER_TXT,YH,True)])
subs=[
    "z ~ 𝒩(0,I) → Flow ODE (T步) → μᵝ 多模态先验",
    "一步策略 μ(s,z) → 动作 a　·　蒸馏 α‖μ−μᵝ‖²",
    "Critic Q_φ(s,a)：TD 学习 → 反馈 max Q",
]
sy=RY+0.42; sh_=0.46; sg=0.07
for i,t in enumerate(subs):
    yy=sy+i*(sh_+sg)
    box(m3x+0.14,yy,m3w-0.28,sh_,SUB_FILL,AMBER,dash=False,lw=0.75,radius=0.06)
    text(m3x+0.14,yy,m3w-0.28,sh_,[(t,9,NAVY,KT,False)])

card(m4x,RY,m4w,RH,"退化感知奖励 r","信息增益 (D-opt/MI)\n− λ₁·可定位性风险(η)\n− λ₂·运动/能耗代价",AMBER,AMBER_TXT)
card(m5x,RY,m5w,RH,"探索决策输出 a*","目标点 / 速度指令",AMBER,AMBER_TXT)

# FQL 行内右向箭头
for (lx) in [m1x+m1w, m2x+m2w, m3x+m3w, m4x+m4w]:
    barrow(lx+0.005,RY+RH/2-0.15,0.11,0.3,MSO_SHAPE.RIGHT_ARROW,AMBER_BAR)

# ===================== 内层框：LIO-SLAM 被动建图 =====================
IX,IY,IW,IH = 0.95,3.62,11.45,3.0
box(IX,IY,IW,IH,IN_FILL,TEAL,dash=False,lw=1.5,radius=0.04)
box(IX,IY,IW,0.38,TEAL,None,radius=0.03)
text(IX+0.18,IY,IW-0.36,0.38,[("内层：LIO-SLAM 建图（被动状态估计）",13,WHITE,YH,True)],align=PP_ALIGN.LEFT)

# 传感器照片（复用）
place(robot,1.15,4.30,1.25,1.16)
place(imu,1.35,5.58,0.9,0.78)
text(1.0,5.0,1.4,0.0,[("",1,GREY,KT,False)])  # spacer (no-op)
# 感知输入 箭头
barrow(2.5,4.95,0.5,0.55,MSO_SHAPE.RIGHT_ARROW,NAVY)
text(2.42,4.62,0.72,0.3,[("感知输入",9,TEAL_D,KT,True)])

# SLAM 流水线
P=[
    ("数据预处理","点云去噪·时间同步·畸变补偿"),
    ("前端状态估计","iEKF 紧耦合·Faster-LIO 直接配准"),
    ("增量地图管理","ikd-Tree · kNN 查询"),
    ("回环 & 位姿图优化","回环检测 · 全局 BA/PGO"),
    ("退化感知与全局输出","Hessian→6-DoF η · 地图/位姿"),
]
px0=3.35; pw=1.42; pg=0.13; pyy=4.45; ph=1.45
for i,(t,sub) in enumerate(P):
    x=px0+i*(pw+pg)
    card(x,pyy,pw,ph,t,sub,TEAL,TEAL_D,fill=WHITE,tsz=10.5,ssz=8.5)
    if i<len(P)-1:
        barrow(x+pw+0.005,pyy+ph/2-0.13,0.11,0.26,MSO_SHAPE.RIGHT_ARROW,NAVY)
# 点云图（复用）放在 P5 右侧
p5_right = px0+4*(pw+pg)+pw
place(pc, p5_right+0.12, 4.85, 0.95, 0.9)
text(p5_right+0.05,5.78,1.1,0.26,[("全局点云地图",8,GREY,KT,True)])

# ===================== 连接：内层 SLAM → 外层决策（生成输入，向上）=====================
barrow(3.45,3.36,0.24,0.24,MSO_SHAPE.UP_ARROW,TEAL)
text(3.72,3.34,4.6,0.28,[("生成输入：SLAM 地图 / 位姿 / η  →  状态编码 sₜ",9.5,TEAL_D,KT,True)],align=PP_ALIGN.LEFT)

# ===================== 闭环反馈：a* → 导航 → 环境 → 新观测（绕内层周边）=====================
lp=AMBER
dline(m5x+m5w/2, RY+RH, 12.55, RY+RH, lp)          # 右行到右边距
dline(12.55, RY+RH, 12.55, 6.92, lp)               # 下行(内层右侧外)
dline(12.55, 6.92, 1.45, 6.92, lp)                 # 左行(内层下方带)
dline(1.45, 6.92, 1.45, 5.5, lp)                   # 上行回到传感器
tri=slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(1.32),Inches(5.46),Inches(0.26),Inches(0.2))
tri.fill.solid(); tri.fill.fore_color.rgb=rgb(lp); tri.line.fill.background(); tri.shadow.inherit=False
text(3.0,6.62,8.0,0.28,[("闭环反馈：a* → 自适应导航执行 → 环境交互 → 新观测",10,AMBER_TXT,YH,True)])

# 把复用的照片提到最上层（否则被内层框填充覆盖）
spTree = slide.shapes._spTree
for p in [robot, imu, pc]:
    if p is not None:
        el = p._element; spTree.remove(el); spTree.append(el)

prs.save(DST)
print("WROTE:",DST)
