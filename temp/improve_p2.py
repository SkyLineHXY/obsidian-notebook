# -*- coding: utf-8 -*-
"""在 研究内容三_总体框架图.pptx 第2页右侧增量添加：退化感知η桥接 + FQL上层决策模块 +
页面标题 + 闭环反馈；并修正 SLAM 前端措辞。另存为 _改进版.pptx，保留原文件与所有图片。"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import copy

SRC = "研究内容三_总体框架图.pptx"
DST = "研究内容三_总体框架图_改进版.pptx"

# ---- palette ----
NAVY   = "0E2841"; TEAL = "1C7293"; TEAL_D = "0F4C5C"; GREY = "5B6B7C"
AMBER  = "E9A23B"; AMBER_BAR = "C77F1A"; AMBER_TXT = "8F5E10"; AMBER_FILL = "FFF6E8"
WHITE  = "FFFFFF"
YH = "Microsoft YaHei"; KT = "楷体"

prs = Presentation(SRC)
slide = prs.slides[1]

def rgb(h): return RGBColor.from_string(h)

def set_cjk(run, name):
    """确保中文 ea/cs typeface 不回退。"""
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)

def style_runs(tf, name):
    for p in tf.paragraphs:
        for r in p.runs:
            set_cjk(r, name)

def box(x, y, w, h, fill, line, dash=False, lw=1.0, radius=0.06):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = rgb(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = rgb(line); sp.line.width = Pt(lw)
        if dash:
            ln = sp.line._get_or_add_ln()
            d = ln.find(qn('a:prstDash'))
            if d is None:
                d = ln.makeelement(qn('a:prstDash'), {}); ln.append(d)
            d.set('val', 'dash')
    sp.shadow.inherit = False
    return sp

def label(x, y, w, h, text, size, color, font=YH, bold=True, align=PP_ALIGN.CENTER,
          anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(1)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = rgb(color)
    set_cjk(r, font)
    return tb

def inner(x, y, w, h, title, sub, line, title_color):
    """白底内框：YaHei 标题 + 楷体副标题。"""
    box(x, y, w, h, WHITE, line, dash=False, lw=1.0, radius=0.08)
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(3); tf.margin_top = Pt(3); tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run(); r1.text = title
    r1.font.size = Pt(11); r1.font.bold = True; r1.font.color.rgb = rgb(title_color); set_cjk(r1, YH)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = sub
    r2.font.size = Pt(9); r2.font.bold = False; r2.font.color.rgb = rgb(GREY); set_cjk(r2, KT)
    return tb

def block_arrow(x, y, w, h, shape=MSO_SHAPE.RIGHT_ARROW, fill=NAVY):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = rgb(fill); sp.line.fill.background()
    sp.shadow.inherit = False
    return sp

def dline(x1, y1, x2, y2, color=AMBER, w=2.0, dash=True):
    cn = slide.shapes.add_connector(2, Inches(x1), Inches(y1), Inches(x2), Inches(y2))  # 2 = straight
    cn.line.color.rgb = rgb(color); cn.line.width = Pt(w)
    if dash:
        ln = cn.line._get_or_add_ln()
        d = ln.makeelement(qn('a:prstDash'), {}); d.set('val', 'dash'); ln.append(d)
    cn.shadow.inherit = False
    return cn

# ============================================================
# 2. 修正 SLAM 前端措辞 (id42)
# ============================================================
for sh in slide.shapes:
    if sh.shape_id == 42 and sh.has_text_frame:
        tf = sh.text_frame
        # 收集原字体（楷体 10pt grey）后整体重写为单段
        # 清空多余 run，保留首段
        p = tf.paragraphs[0]
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
        # 删除额外段落
        for extra in tf.paragraphs[1:]:
            extra._p.getparent().remove(extra._p)
        nr = p.add_run(); nr.text = "iEKF 紧耦合 · Faster-LIO 直接配准"
        nr.font.size = Pt(10); nr.font.bold = True; nr.font.color.rgb = rgb(GREY); set_cjk(nr, KT)
        break

# ============================================================
# 1. 页面标题
# ============================================================
label(0.4, 0.22, 12.2, 0.55, "图2  总体方法框架：退化感知驱动的具身主动 SLAM 闭环",
      18, NAVY, font=YH, bold=True, align=PP_ALIGN.LEFT)

# ============================================================
# 3. 退化感知 η 桥接框（琥珀）
# ============================================================
eta_x, eta_y, eta_w, eta_h = 7.15, 3.55, 1.95, 1.45
box(eta_x, eta_y, eta_w, eta_h, AMBER_FILL, AMBER, dash=True, lw=1.25, radius=0.08)
tb = slide.shapes.add_textbox(Inches(eta_x), Inches(eta_y), Inches(eta_w), Inches(eta_h))
tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = tf.margin_right = Pt(4)
p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
r1 = p1.add_run(); r1.text = "退化感知量化"
r1.font.size = Pt(11); r1.font.bold = True; r1.font.color.rgb = rgb(AMBER_TXT); set_cjk(r1, YH)
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = "Hessian 特征空间分解\n→ 6-DoF 可定位性 η\n(none / partial / full)"
r2.font.size = Pt(9); r2.font.color.rgb = rgb(GREY); set_cjk(r2, KT)

# 入箭头：SLAM 输出 → η 框
block_arrow(6.62, 4.18, 0.46, 0.62, MSO_SHAPE.RIGHT_ARROW)
label(6.35, 3.82, 1.0, 0.3, "全局地图/位姿", 8.5, TEAL_D, font=KT, bold=True)

# ============================================================
# 4. FQL 上层决策分组框（琥珀）
# ============================================================
fx, fy, fw, fh = 9.35, 1.95, 3.55, 4.25
box(fx, fy, fw, fh, AMBER_FILL, AMBER, dash=True, lw=1.25, radius=0.06)
# 标题条
bar = box(fx, fy, fw, 0.36, AMBER_BAR, None, dash=False, radius=0.04)
label(fx + 0.08, fy, fw - 0.16, 0.36, "主动感知与决策（FQL 上层）", 13, WHITE, font=YH, bold=True)

# 4 内框
items = [
    ("感知状态构建 sₜ", "地图熵 H(m)·协方差 Σ·可定位性 η·候选(前沿+回环)"),
    ("FQL 双策略决策", "BC Flow 多模态先验 ⊕ 一步策略 max Q · 无 BPTT"),
    ("退化感知奖励 r", "信息增益 − λ₁·可定位性风险 − λ₂·运动代价"),
    ("探索决策输出 a*", "目标点·速度 → 自适应导航执行"),
]
ix = fx + 0.18; iw = fw - 0.36
top = fy + 0.5; bh = 0.78; gap = 0.18
ys = []
for i, (t, sub) in enumerate(items):
    yy = top + i * (bh + gap)
    ys.append(yy)
    inner(ix, yy, iw, bh, t, sub, AMBER, AMBER_TXT)
# 向下小块箭头
for i in range(3):
    ay = ys[i] + bh + 0.005
    block_arrow(fx + fw / 2 - 0.11, ay, 0.22, gap - 0.02, MSO_SHAPE.DOWN_ARROW)

# 入箭头（生成输入）：η 框 → FQL
block_arrow(eta_x + eta_w + 0.02, 4.0, 0.46, 0.6, MSO_SHAPE.RIGHT_ARROW)
label(eta_x + eta_w - 0.02, 3.62, 0.62, 0.3, "生成输入", 8.5, AMBER_TXT, font=KT, bold=True)

# ============================================================
# 5. 闭环反馈（琥珀虚线，FQL顶 → 顶部 → 传感器）
# ============================================================
loop_y = 1.05
dline(fx + 1.6, fy, fx + 1.6, loop_y)            # 上行
dline(fx + 1.6, loop_y, 1.4, loop_y)             # 左行
dline(1.4, loop_y, 1.4, 2.62)                    # 下行（到机器人照片上方）
# 末端箭头（向下等腰三角）
tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(1.27), Inches(2.6), Inches(0.26), Inches(0.2))
tri.rotation = 180; tri.fill.solid(); tri.fill.fore_color.rgb = rgb(AMBER); tri.line.fill.background(); tri.shadow.inherit = False
label(4.2, loop_y - 0.34, 6.2, 0.3, "闭环：主动探索 a* → 环境交互 → 新观测", 10, AMBER_TXT, font=YH, bold=True, align=PP_ALIGN.CENTER)

prs.save(DST)
print("WROTE:", DST)
